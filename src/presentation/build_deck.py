"""Generate the project presentation (presentation.pptx) from pipeline outputs.

Builds a 7-slide academic-style deck mirroring the structure of the project
report (README.md): introduction & task, dataset & features, solution
architecture, automation, results, monitoring & reproducibility, conclusions.

Reads reports/metrics.json + reports/figures/*.png. Run after the pipeline:

    python -m src.presentation.build_deck
"""

from __future__ import annotations

import json

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from ..utils import get_logger, resolve

LOG = get_logger()

NAVY = RGBColor(0x1F, 0x2D, 0x5A)
BLUE = RGBColor(0x4C, 0x72, 0xB0)
GREY = RGBColor(0x55, 0x55, 0x55)
DARK = RGBColor(0x22, 0x22, 0x22)


def _load_metrics() -> dict:
    path = resolve("reports/metrics.json")
    if not path.exists():
        raise FileNotFoundError("reports/metrics.json not found — run the pipeline first.")
    with open(path, encoding="utf-8") as fh:
        return json.load(fh)


def _add_title(slide, text: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.1), Inches(0.9))
    run = box.text_frame.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(30)
    run.font.bold = True
    run.font.color.rgb = NAVY


def _add_textbox(slide, left: float, top: float, width: float, height: float,
                  paragraphs: list[tuple[str, int, bool, RGBColor]]) -> None:
    """paragraphs is a list of (text, font_size_pt, bold, color)."""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, (text, size, bold, color) in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        p.space_after = Pt(6)


def _add_bullets(slide, left: float, top: float, width: float, height: float,
                  bullets: list[str], size: int = 16) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = "•  " + text
        run.font.size = Pt(size)
        run.font.color.rgb = DARK
        p.space_after = Pt(6)


def _add_image(slide, name: str, left: float, top: float, width: float) -> None:
    path = resolve(f"reports/figures/{name}")
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width))


def _title_slide(prs: Presentation) -> None:
    """Слайд 1. Титульный."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Title
    box = slide.shapes.add_textbox(Inches(0.7), Inches(2.0), Inches(12), Inches(1.3))
    run = box.text_frame.paragraphs[0].add_run()
    run.text = "Автоматизация ML-пайплайна"
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = NAVY
    # Subtitle
    box2 = slide.shapes.add_textbox(Inches(0.7), Inches(3.3), Inches(12), Inches(0.8))
    run2 = box2.text_frame.paragraphs[0].add_run()
    run2.text = "Knowledge tracing на датасете ASSISTments 2009"
    run2.font.size = Pt(24)
    run2.font.color.rgb = BLUE
    # Author + date
    box3 = slide.shapes.add_textbox(Inches(0.7), Inches(5.5), Inches(12), Inches(1.5))
    tf = box3.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = "Автор: Григорьева Е.С."
    r1.font.size = Pt(18)
    r1.font.color.rgb = GREY
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = "Май 2026"
    r2.font.size = Pt(18)
    r2.font.color.rgb = GREY


def _task_slide(prs: Presentation, metrics: dict) -> None:
    """Слайд 2. Постановка задачи."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "Постановка задачи")
    _add_textbox(slide, 0.6, 1.3, 12.1, 1.0, [
        ("Предметная область", 17, True, NAVY),
    ])
    _add_textbox(slide, 0.6, 1.7, 12.1, 1.4, [
        ("Knowledge tracing — задача образовательной аналитики: по истории",
         15, False, DARK),
        ("взаимодействий студента (тройки «студент, навык, верно/неверно»),",
         15, False, DARK),
        ("упорядоченных по времени, оценить уровень владения каждым навыком",
         15, False, DARK),
        ("и предсказать результат следующего ответа (one-step-ahead).",
         15, False, DARK),
    ])
    _add_textbox(slide, 0.6, 3.5, 12.1, 0.5, [
        ("Формализация ML-задачи", 17, True, NAVY),
    ])
    _add_bullets(slide, 0.6, 4.0, 12.1, 2.7, [
        "Тип: бинарная классификация на временных рядах (correct ∈ {0, 1}).",
        "Вход: история до момента t — последовательность (skill_id, correct) или табличные причинные признаки.",
        "Выход: вероятность правильного ответа на шаге t + 1.",
        "Основная метрика: ROC-AUC (устойчива к несбалансированности классов).",
        "Дополнительно: Accuracy, F1, Precision, Recall, RMSE, log-loss.",
        "Сплит выборки выполняется по студентам, что исключает утечку между train, val и test.",
    ], size=15)


def _data_slide(prs: Presentation, metrics: dict) -> None:
    """Слайд 3. Данные и признаки."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "Данные и признаки")
    ds = metrics["dataset"]
    _add_textbox(slide, 0.6, 1.3, 6.4, 0.5, [
        ("Датасет ASSISTments 2009 (skill-builder)", 17, True, NAVY),
    ])
    _add_bullets(slide, 0.6, 1.85, 6.4, 2.3, [
        f"{ds['n_students']} студентов · {ds['n_interactions']} взаимодействий · {ds['n_skills']} навыков.",
        f"Глобальная доля верных ответов: {ds['global_correct_rate']:.4f}.",
        f"Сплит по студентам 70/10/20: train {ds['split_students']['train']}, "
        f"val {ds['split_students']['val']}, test {ds['split_students']['test']}.",
        "Закоммичен sample на 288 студентов для офлайн-тестов и CI.",
    ], size=14)
    _add_textbox(slide, 0.6, 4.2, 6.4, 0.5, [
        ("Восемь причинных признаков", 17, True, NAVY),
    ])
    _add_bullets(slide, 0.6, 4.75, 6.4, 2.5, [
        "skill_idx, order_idx — кто и когда.",
        "user_prior_attempts, user_prior_correct_rate — история студента.",
        "skill_prior_attempts, skill_prior_correct_rate — история по навыку.",
        "recent3_correct_rate — успешность за 3 последних шага.",
        "skill_difficulty (вычислена только на train, чтобы избежать утечки).",
    ], size=13)
    _add_image(slide, "dataset_overview.png", left=7.3, top=1.5, width=5.7)


def _architecture_slide(prs: Presentation) -> None:
    """Слайд 4. Архитектура решения + автоматизация."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "Архитектура решения и три уровня автоматизации")
    _add_textbox(slide, 0.6, 1.3, 12.1, 0.5, [
        ("Поток данных", 17, True, NAVY),
    ])
    _add_textbox(slide, 0.6, 1.8, 12.1, 0.5, [
        ("Extract → Transform → Load → Data-quality gate → 4 модели → Оценка → Визуализация → MLflow",
         13, False, DARK),
    ])
    _add_textbox(slide, 0.6, 2.6, 12.1, 0.5, [
        ("Четыре модели разной природы", 17, True, NAVY),
    ])
    _add_bullets(slide, 0.6, 3.1, 12.1, 1.8, [
        "BKT — скрытая марковская модель (EM-обучение, один набор параметров на навык). Интерпретируемый baseline.",
        "DKT — рекуррентная нейросеть (Embedding → LSTM → Linear). Моделирует длинные зависимости.",
        "DKT + Optuna — TPE-поиск архитектуры (embed_dim, hidden_dim, dropout, lr, batch_size).",
        "AutoML (FLAML) — выбор модели и гиперпараметров за фиксированный бюджет времени.",
    ], size=14)
    _add_textbox(slide, 0.6, 5.1, 12.1, 0.5, [
        ("Три уровня автоматизации", 17, True, NAVY),
    ])
    _add_bullets(slide, 0.6, 5.6, 12.1, 1.6, [
        "1. Выбор модели и гиперпараметров — FLAML на табличных признаках.",
        "2. Поиск архитектуры нейросети — Optuna поверх DKT (25 trials).",
        "3. Сквозная оркестрация — один запуск pipeline.py выполняет весь цикл.",
    ], size=14)


def _results_slide(prs: Presentation, metrics: dict) -> None:
    """Слайд 5. Результаты экспериментов."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "Результаты на отложенной тестовой выборке")
    results = metrics["results"]
    best = metrics["best_model"]

    # Table
    rows = len(results) + 1
    table = slide.shapes.add_table(rows, 5, Inches(0.6), Inches(1.4), Inches(6.5),
                                    Inches(0.4 * rows)).table
    headers = ["Модель", "AUC", "Accuracy", "F1", "RMSE"]
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.bold = True
                r.font.size = Pt(13)
                r.font.color.rgb = NAVY
    for i, (name, m) in enumerate(results.items(), start=1):
        bold = (name == best)
        vals = [name, f"{m['auc']:.4f}", f"{m['accuracy']:.4f}", f"{m['f1']:.4f}", f"{m['rmse']:.4f}"]
        for j, v in enumerate(vals):
            cell = table.cell(i, j)
            cell.text = v
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(12)
                    r.font.bold = bold
                    r.font.color.rgb = NAVY if bold else DARK

    _add_image(slide, "model_comparison.png", left=7.3, top=1.4, width=5.7)

    # Interpretation
    _add_textbox(slide, 0.6, 5.0, 12.1, 0.5, [
        ("Интерпретация", 16, True, NAVY),
    ])
    _add_bullets(slide, 0.6, 5.45, 12.1, 1.8, [
        f"Лучший результат — {best} (AUC {results[best]['auc']:.4f}); соответствует "
        f"уровню публикуемых бенчмарков для DKT на ASSISTments 2009.",
        f"AutoML выбрал {metrics['automl_best_estimator']}; уступает DKT на 0.010 AUC "
        f"при полностью автоматической настройке.",
        "Optuna нашёл конфигурацию с сильной регуляризацией (dropout 0.49), "
        "что указывает на склонность DKT к переобучению на данном объёме данных.",
    ], size=13)


def _monitoring_slide(prs: Presentation, metrics: dict) -> None:
    """Слайд 6. Мониторинг и воспроизводимость."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "Мониторинг и воспроизводимость")
    drift = metrics["drift"]
    res = metrics["resources"]

    _add_textbox(slide, 0.6, 1.3, 6.4, 0.5, [
        ("Мониторинг качества данных", 16, True, NAVY),
    ])
    _add_bullets(slide, 0.6, 1.8, 6.4, 1.8, [
        "Data-quality gate: 5 проверок (схема, пропуски, бинарные метки, диапазоны, дубликаты).",
        f"Статус текущего прогона: passed = {metrics['data_quality_passed']}.",
        f"Мониторинг дрейфа: PSI + KS-тест на {drift['n_features']} признаках.",
        f"Дрейфующих признаков: {drift['n_significant_drift']} ({drift['overall_status']}).",
    ], size=13)

    _add_textbox(slide, 0.6, 3.8, 6.4, 0.5, [
        ("Мониторинг ресурсов (psutil)", 16, True, NAVY),
    ])
    _add_bullets(slide, 0.6, 4.3, 6.4, 2.3, [
        f"BKT: {res['BKT']['elapsed_s']:.0f} с · {res['BKT']['peak_rss_mb']:.0f} MB · "
        f"{res['BKT']['avg_cpu_pct']:.0f}% CPU.",
        f"DKT: {res['DKT']['elapsed_s']:.0f} с · {res['DKT']['peak_rss_mb']:.0f} MB · "
        f"{res['DKT']['avg_cpu_pct']:.0f}% CPU.",
        f"DKT + Optuna: {res['DKT+Optuna']['elapsed_s']:.0f} с · "
        f"{res['DKT+Optuna']['peak_rss_mb']:.0f} MB · {res['DKT+Optuna']['avg_cpu_pct']:.0f}% CPU.",
        f"AutoML (FLAML): {res['AutoML']['elapsed_s']:.0f} с · "
        f"{res['AutoML']['peak_rss_mb']:.0f} MB · {res['AutoML']['avg_cpu_pct']:.0f}% CPU.",
    ], size=13)

    _add_textbox(slide, 7.1, 1.3, 5.9, 0.5, [
        ("Воспроизводимость", 16, True, NAVY),
    ])
    _add_bullets(slide, 7.1, 1.8, 5.9, 2.0, [
        "Docker: образ python:3.11-slim, non-root user, CPU-сборка torch.",
        "docker-compose: сервисы pipeline + MLflow UI на одном томе.",
        "CI/CD GitHub Actions: lint (ruff) + pytest + docker build + smoke-прогон.",
        "MLflow: трекинг параметров, метрик, артефактов и обученных моделей.",
    ], size=13)

    _add_textbox(slide, 7.1, 4.0, 5.9, 0.5, [
        ("Тестирование", 16, True, NAVY),
    ])
    _add_bullets(slide, 7.1, 4.5, 5.9, 2.0, [
        "7 файлов на pytest, маркер slow для тяжёлых тестов.",
        "Покрытие: ETL, Transform, метрики, BKT, мониторинг, smoke DKT/FLAML.",
        "Линтер ruff (правила E, F, I, W).",
    ], size=13)


def _conclusions_slide(prs: Presentation, metrics: dict) -> None:
    """Слайд 7. Выводы."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "Выводы")
    best = metrics["best_model"]
    best_auc = metrics["results"][best]["auc"]

    _add_textbox(slide, 0.6, 1.3, 12.1, 0.5, [
        ("Технические итоги", 17, True, NAVY),
    ])
    _add_bullets(slide, 0.6, 1.8, 12.1, 2.2, [
        "Реализован сквозной автоматизированный ML-пайплайн: ETL, четыре модели, оценка, мониторинг, упаковка.",
        f"Лучший результат — {best} с AUC = {best_auc:.4f} на отложенной тестовой выборке.",
        "Автоматизация на трёх уровнях: FLAML (модель + гиперпараметры), Optuna (архитектура DKT), pipeline.py (оркестрация).",
        "Мониторинг качества данных, дрейфа (PSI + KS), ресурсов; всё логируется в MLflow.",
        "Воспроизводимость обеспечена контейнером Docker и CI на GitHub Actions.",
    ], size=14)

    _add_textbox(slide, 0.6, 4.4, 12.1, 0.5, [
        ("Содержательные итоги", 17, True, NAVY),
    ])
    _add_bullets(slide, 0.6, 4.9, 12.1, 2.4, [
        "Рекуррентные модели (DKT) превосходят BKT на 0.08–0.09 AUC; разрыв соответствует литературе.",
        "Автоматический подбор архитектуры (Optuna) даёт прирост +0.002 AUC и выбирает сильную регуляризацию.",
        "AutoML на инженерных причинных признаках достигает AUC 0.7914 — на 0.010 ниже специализированной DKT.",
        "Дрейфа между train и test не обнаружено; data-quality gate пройден.",
    ], size=14)


def build(output: str = "presentation/presentation.pptx") -> str:
    metrics = _load_metrics()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    _title_slide(prs)
    _task_slide(prs, metrics)
    _data_slide(prs, metrics)
    _architecture_slide(prs)
    _results_slide(prs, metrics)
    _monitoring_slide(prs, metrics)
    _conclusions_slide(prs, metrics)

    out = resolve(output)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    LOG.info("Presentation saved -> %s (%d slides)", out, len(prs.slides._sldIdLst))
    return str(out)


if __name__ == "__main__":
    build()
